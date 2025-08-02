from flask import Flask, request, jsonify, send_file, after_this_request
from flask_cors import CORS
import requests
import io
import zipfile
import os
import tempfile
import gc
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import logging
import threading
import time

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app)

# Global list to track temporary files and buffers for cleanup
cleanup_queue = []
cleanup_lock = threading.Lock()

def cleanup_resources():
    """Background cleanup of temporary resources"""
    while True:
        try:
            with cleanup_lock:
                current_time = time.time()
                items_to_remove = []
                
                for item in cleanup_queue:
                    # Clean up items older than 5 minutes
                    if current_time - item['created'] > 300:
                        try:
                            if item['type'] == 'file' and os.path.exists(item['path']):
                                os.unlink(item['path'])
                                logger.info(f"Cleaned up temporary file: {item['path']}")
                            elif item['type'] == 'buffer':
                                if hasattr(item['buffer'], 'close'):
                                    item['buffer'].close()
                                logger.info("Cleaned up buffer")
                        except Exception as e:
                            logger.error(f"Error cleaning up {item['type']}: {e}")
                        finally:
                            items_to_remove.append(item)
                
                # Remove cleaned items from queue
                for item in items_to_remove:
                    cleanup_queue.remove(item)
            
            # Force garbage collection
            if items_to_remove:
                gc.collect()
                
            time.sleep(60)  # Check every minute
            
        except Exception as e:
            logger.error(f"Cleanup thread error: {e}")
            time.sleep(60)

# Start cleanup thread
cleanup_thread = threading.Thread(target=cleanup_resources, daemon=True)
cleanup_thread.start()

def schedule_cleanup(resource_type, path=None, buffer=None):
    """Schedule a resource for cleanup"""
    with cleanup_lock:
        cleanup_queue.append({
            'type': resource_type,
            'path': path,
            'buffer': buffer,
            'created': time.time()
        })

def download_image(url):
    """Download image from URL and return PIL Image object"""
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        response = requests.get(url, headers=headers, timeout=30)
        response.raise_for_status()
        
        # Convert to PIL Image
        image = Image.open(io.BytesIO(response.content))
        return image, response.content
    except Exception as e:
        logger.error(f"Failed to download image {url}: {str(e)}")
        raise

def create_pdf_from_images(images, title):
    """Create PDF from list of image URLs"""
    pdf_buffer = io.BytesIO()
    
    try:
        # Download first image to get dimensions
        first_image, _ = download_image(images[0])
        width, height = first_image.size
        
        # Create PDF with custom page size
        c = canvas.Canvas(pdf_buffer, pagesize=(width, height))
        
        for i, image_url in enumerate(images):
            logger.info(f"Processing image {i+1}/{len(images)}")
            
            try:
                image, image_data = download_image(image_url)
                
                # Save image to temporary file
                with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_file:
                    temp_path = temp_file.name
                    image.save(temp_path, 'JPEG')
                    
                    # Add image to PDF
                    c.drawImage(temp_path, 0, 0, width=width, height=height)
                    
                    # Schedule cleanup
                    schedule_cleanup('file', path=temp_path)
                
                if i < len(images) - 1:  # Don't add page after last image
                    c.showPage()
                    
            except Exception as e:
                logger.error(f"Failed to process image {i+1}: {str(e)}")
                continue
        
        c.save()
        pdf_buffer.seek(0)
        return pdf_buffer
        
    except Exception as e:
        logger.error(f"PDF creation failed: {str(e)}")
        raise

def create_ppt_from_images(images, title):
    """Create PowerPoint from list of image URLs"""
    try:
        prs = Presentation()
        
        # Remove default slide layout
        prs.slide_layouts._sldLayouts.clear()
        
        for i, image_url in enumerate(images):
            logger.info(f"Processing slide {i+1}/{len(images)}")
            
            try:
                image, image_data = download_image(image_url)
                width, height = image.size
                
                # Create slide with custom size
                slide_layout = prs.slide_layouts[0] if prs.slide_layouts else None
                if slide_layout is None:
                    # Create blank slide
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                else:
                    slide = prs.slides.add_slide(slide_layout)
                
                # Clear existing shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'element'):
                        shape.element.getparent().remove(shape.element)
                
                # Save image to temporary file
                with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_file:
                    temp_path = temp_file.name
                    image.save(temp_path, 'JPEG')
                    
                    # Add image to slide (full slide)
                    slide.shapes.add_picture(
                        temp_path, 
                        Inches(0), 
                        Inches(0),
                        width=Inches(10),  # Standard slide width
                        height=Inches(7.5)  # Standard slide height
                    )
                    
                    # Schedule cleanup
                    schedule_cleanup('file', path=temp_path)
                    
            except Exception as e:
                logger.error(f"Failed to process slide {i+1}: {str(e)}")
                continue
        
        # Save to buffer
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer
        
    except Exception as e:
        logger.error(f"PPT creation failed: {str(e)}")
        raise

def create_zip_from_images(images, title):
    """Create ZIP file from list of image URLs"""
    try:
        zip_buffer = io.BytesIO()
        
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            for i, image_url in enumerate(images):
                logger.info(f"Downloading image {i+1}/{len(images)}")
                
                try:
                    image, image_data = download_image(image_url)
                    
                    # Save image with proper filename
                    filename = f"slide_{str(i+1).zfill(3)}.jpg"
                    
                    # Convert to JPEG if needed
                    if image.mode in ('RGBA', 'LA'):
                        # Convert RGBA to RGB
                        rgb_image = Image.new('RGB', image.size, (255, 255, 255))
                        rgb_image.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
                        image = rgb_image
                    
                    # Save to buffer
                    img_buffer = io.BytesIO()
                    image.save(img_buffer, 'JPEG', quality=95)
                    img_buffer.seek(0)
                    
                    # Add to ZIP
                    zip_file.writestr(filename, img_buffer.getvalue())
                    
                except Exception as e:
                    logger.error(f"Failed to add image {i+1} to ZIP: {str(e)}")
                    continue
        
        zip_buffer.seek(0)
        return zip_buffer
        
    except Exception as e:
        logger.error(f"ZIP creation failed: {str(e)}")
        raise

@app.route('/convert', methods=['POST'])
def convert():
    try:
        data = request.get_json()
        
        if not data:
            return jsonify({'error': 'No JSON data provided'}), 400
        
        images = data.get('images', [])
        title = data.get('title', 'Presentation')
        format_type = data.get('format', 'pdf').lower()
        
        if not images:
            return jsonify({'error': 'No images provided'}), 400
        
        if format_type not in ['pdf', 'ppt', 'zip']:
            return jsonify({'error': 'Invalid format. Must be pdf, ppt, or zip'}), 400
        
        logger.info(f"Converting {len(images)} images to {format_type}")
        
        # Create appropriate file
        file_buffer = None
        if format_type == 'pdf':
            file_buffer = create_pdf_from_images(images, title)
            mimetype = 'application/pdf'
            filename = f"{title}.pdf"
            
        elif format_type == 'ppt':
            file_buffer = create_ppt_from_images(images, title)
            mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            filename = f"{title}.pptx"
            
        elif format_type == 'zip':
            file_buffer = create_zip_from_images(images, title)
            mimetype = 'application/zip'
            filename = f"{title}.zip"
        
        # Schedule buffer cleanup after response
        if file_buffer:
            schedule_cleanup('buffer', buffer=file_buffer)
            
            @after_this_request
            def cleanup_after_response(response):
                """Cleanup resources after response is sent"""
                try:
                    # Force garbage collection
                    gc.collect()
                    logger.info(f"Response sent, scheduled cleanup for {filename}")
                except Exception as e:
                    logger.error(f"Post-response cleanup error: {e}")
                return response
        
        # Return file
        return send_file(
            file_buffer,
            mimetype=mimetype,
            as_attachment=True,
            download_name=filename
        )
        
    except Exception as e:
        logger.error(f"Conversion error: {str(e)}")
        return jsonify({'error': f'Conversion failed: {str(e)}'}), 500

@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'healthy', 'service': 'slidesdown-converter'})

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=False)
